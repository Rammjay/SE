import os
from flask import Flask, request, jsonify
from flask_cors import CORS
import openai
import sqlite3
import traceback
from werkzeug.utils import secure_filename
import docx
import PyPDF2
import pptx
from transformers import pipeline

from dotenv import load_dotenv
load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")
print("OpenAI API Key:", openai.api_key)  # <-- add this here temporarily

app = Flask(__name__)
CORS(app)
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

openai.api_key = os.getenv("OPENAI_API_KEY")

# Initialize the summarization pipeline
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

def get_db_connection():
    conn = sqlite3.connect('campus.db')
    conn.row_factory = sqlite3.Row
    return conn

def fetch_schedule():
    conn = get_db_connection()
    schedule = conn.execute('SELECT * FROM schedule').fetchall()
    conn.close()
    return schedule

def fetch_events():
    conn = get_db_connection()
    events = conn.execute('SELECT * FROM events').fetchall()
    conn.close()
    return events

def generate_gpt_response(user_input, schedule_info, event_info):
    # Compose prompt for GPT
    prompt = f"""
You are a helpful campus assistant.

User asked: "{user_input}"

Here are the current class schedules:
{schedule_info}

Here are upcoming events:
{event_info}

Please respond naturally and helpfully.
"""

    response = openai.Completion.create(
        engine="text-davinci-003",  # Or the latest model you want to use
        prompt=prompt,
        max_tokens=150,
        temperature=0.7,
    )

    answer = response.choices[0].text.strip()
    return answer

@app.route('/process-voice', methods=['POST'])
def process_voice():
    try:
        data = request.get_json()
        user_input = data.get('text', '')

        # Fetch schedule and event info from DB
        schedule = fetch_schedule()
        events = fetch_events()

        # Format schedule info nicely for prompt
        schedule_info = "\n".join(
            [f"{row['class_name']} at {row['start_time']} in {row['location']}" for row in schedule]
        )
        event_info = "\n".join(
            [f"{row['event_name']} at {row['event_time']} in {row['location']}" for row in events]
        )

        # Get GPT generated response
        reply = generate_gpt_response(user_input, schedule_info, event_info)

        return jsonify({'reply': reply})
    
    except Exception as e:
        import traceback
        print("=== ERROR OCCURRED ===")
        print(e)
        traceback.print_exc()
        print("======================")
        return jsonify({'error': str(e)}), 500

def extract_text_from_file(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_extension == '.pdf':
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ''
                for page in pdf_reader.pages:
                    text += page.extract_text()
                return text
        
        elif file_extension in ['.doc', '.docx']:
            doc = docx.Document(file_path)
            text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            return text
        
        elif file_extension in ['.ppt', '.pptx']:
            prs = pptx.Presentation(file_path)
            text = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
            return text
        
        else:
            return None
    except Exception as e:
        print(f"Error extracting text: {str(e)}")
        return None

@app.route('/summarize', methods=['POST'])
def summarize_document():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        # Extract text from the document
        text = extract_text_from_file(file_path)
        if text is None:
            return jsonify({'error': 'Unsupported file format or extraction failed'}), 400
        
        # Clean and chunk the text (BART has a max input length)
        # Split into chunks of roughly 1000 characters
        chunks = [text[i:i + 1000] for i in range(0, len(text), 1000)]
        summaries = []
        
        # Summarize each chunk
        for chunk in chunks[:3]:  # Limit to first 3 chunks to keep it manageable
            if len(chunk.strip()) > 100:  # Only summarize chunks with substantial content
                summary = summarizer(chunk, max_length=130, min_length=30, do_sample=False)
                summaries.append(summary[0]['summary_text'])
        
        # Combine the summaries
        final_summary = ' '.join(summaries)
        
        # Clean up the uploaded file
        os.remove(file_path)
        
        return jsonify({'summary': final_summary})
    
    except Exception as e:
        print("=== ERROR OCCURRED ===")
        print(e)
        traceback.print_exc()
        print("======================")
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000)
